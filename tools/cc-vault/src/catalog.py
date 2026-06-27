"""
Catalog Scanner - Walk library directories, hash files, extract text, summarize.

Streaming protocol: when stream=True, prints one JSON line per file to stdout:
  {"event":"progress","phase":"scan","processed":47,"total":852,"file":"board-minutes.pdf","status":"new"}
  {"event":"complete","phase":"scan","new":523,"updated":12,"skipped":317,"errors":0}
"""

import hashlib
import json
import os
import sys
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

try:
    from . import db
    from .config import (
        CATALOG_SUMMARIZABLE_EXTENSIONS,
        CATALOG_METADATA_ONLY_EXTENSIONS,
    )
except ImportError:
    import db
    from config import (
        CATALOG_SUMMARIZABLE_EXTENSIONS,
        CATALOG_METADATA_ONLY_EXTENSIONS,
    )

# All extensions we track (summarizable + metadata-only)
ALL_TRACKED_EXTENSIONS = CATALOG_SUMMARIZABLE_EXTENSIONS | CATALOG_METADATA_ONLY_EXTENSIONS


def _hash_file(path: str, chunk_size: int = 8192) -> str:
    """Compute SHA-256 hash of a file."""
    h = hashlib.sha256()
    with open(path, 'rb') as f:
        while True:
            chunk = f.read(chunk_size)
            if not chunk:
                break
            h.update(chunk)
    return h.hexdigest()


def _stream_event(event: dict) -> None:
    """Print a JSON event line to stdout for UI consumption."""
    print(json.dumps(event, ensure_ascii=True), flush=True)


def _derive_department(file_path: str, library_path: str) -> Optional[str]:
    """Derive department from first subfolder relative to library path."""
    try:
        rel = os.path.relpath(file_path, library_path)
        parts = Path(rel).parts
        if len(parts) > 1:
            return parts[0]
    except ValueError:
        pass
    return None


class CatalogScanner:
    """Scan library directories, catalog files, and generate summaries."""

    def scan_library_by_label(self, label: str, stream: bool = False) -> Dict[str, Any]:
        """Scan a library by its label."""
        lib = db.get_library(label)
        if not lib:
            raise ValueError(f"Library '{label}' not found")
        return self.scan_library(lib['id'], stream=stream)

    def scan_library(self, library_id: int, stream: bool = False) -> Dict[str, Any]:
        """Walk directory, hash files, upsert catalog entries.

        If stream=True, print JSON progress lines to stdout for UI consumption.
        """
        lib = db.get_library_by_id(library_id)
        if not lib:
            raise ValueError(f"Library ID {library_id} not found")

        library_path = lib['path']
        recursive = bool(lib.get('recursive', 1))

        if not os.path.isdir(library_path):
            raise ValueError(f"Library path does not exist: {library_path}")

        # Phase 1: collect all trackable files
        file_list = []
        if recursive:
            for root, _dirs, files in os.walk(library_path):
                for fname in files:
                    ext = os.path.splitext(fname)[1].lower()
                    if ext in ALL_TRACKED_EXTENSIONS:
                        file_list.append(os.path.join(root, fname))
        else:
            for fname in os.listdir(library_path):
                full = os.path.join(library_path, fname)
                if os.path.isfile(full):
                    ext = os.path.splitext(fname)[1].lower()
                    if ext in ALL_TRACKED_EXTENSIONS:
                        file_list.append(full)

        total = len(file_list)
        counts = {"new": 0, "updated": 0, "skipped": 0, "errors": 0, "missing": 0}
        existing_paths = set()

        for idx, fpath in enumerate(file_list, 1):
            fname = os.path.basename(fpath)
            ext = os.path.splitext(fname)[1].lower()
            existing_paths.add(fpath)

            try:
                stat = os.stat(fpath)
                file_size = stat.st_size
                file_modified = datetime.fromtimestamp(stat.st_mtime).isoformat()
                file_hash = _hash_file(fpath)

                # Check if already cataloged
                existing = db.get_catalog_entry_by_path(fpath)
                summarizable = ext in CATALOG_SUMMARIZABLE_EXTENSIONS
                department = _derive_department(fpath, library_path)

                if existing and existing.get('file_hash') == file_hash:
                    status_str = "skipped"
                    counts["skipped"] += 1
                elif existing:
                    status_str = "updated"
                    counts["updated"] += 1
                else:
                    status_str = "new"
                    counts["new"] += 1

                entry_status = 'pending' if summarizable else 'skipped'
                # If file hash unchanged and entry was already summarized/skipped, preserve
                if existing and existing.get('file_hash') == file_hash:
                    entry_status = existing.get('status', entry_status)

                db.upsert_catalog_entry(
                    library_id=library_id,
                    file_path=fpath,
                    file_name=fname,
                    file_ext=ext,
                    file_size=file_size,
                    file_hash=file_hash,
                    file_modified_at=file_modified,
                    department=department,
                    summarizable=summarizable,
                    status=entry_status,
                )

                if stream:
                    _stream_event({
                        "event": "progress",
                        "phase": "scan",
                        "processed": idx,
                        "total": total,
                        "file": fname,
                        "status": status_str,
                    })

            except Exception as e:
                counts["errors"] += 1
                if stream:
                    _stream_event({
                        "event": "progress",
                        "phase": "scan",
                        "processed": idx,
                        "total": total,
                        "file": fname,
                        "status": "error",
                        "error": str(e),
                    })

        # Mark missing entries
        missing_count = db.mark_missing_catalog_entries(library_id, existing_paths)
        counts["missing"] = missing_count

        # Update last_scanned
        db.update_library_last_scanned(library_id)

        if stream:
            _stream_event({
                "event": "complete",
                "phase": "scan",
                **counts,
            })

        return counts

    def summarize_entries(self, library_id: Optional[int] = None,
                          batch_size: int = 10, stream: bool = False,
                          dry_run: bool = False) -> Dict[str, Any]:
        """Generate AI summaries for pending catalog entries.

        Processes ALL pending entries in batches of batch_size until none remain.
        """
        if dry_run:
            entries = db.get_pending_catalog_entries(library_id=library_id, limit=batch_size)
            return {"pending": len(entries)}

        counts = {"summarized": 0, "deduped": 0, "errors": 0}
        processed_total = 0

        while True:
            entries = db.get_pending_catalog_entries(library_id=library_id, limit=batch_size)
            if not entries:
                break

            for idx, entry in enumerate(entries, 1):
                processed_total += 1
                try:
                    # Dedup check: same hash already summarized?
                    if entry.get('file_hash'):
                        donor = db.get_catalog_entry_by_hash(entry['file_hash'])
                        if donor and donor['id'] != entry['id']:
                            # Copy summary from donor
                            db.update_catalog_entry_summary(
                                entry_id=entry['id'],
                                title=donor.get('title', ''),
                                summary=donor.get('summary', ''),
                                tags=donor.get('tags', ''),
                                dedup_source_id=donor['id'],
                            )
                            counts["deduped"] += 1
                            if stream:
                                _stream_event({
                                    "event": "progress",
                                    "phase": "summarize",
                                    "processed": processed_total,
                                    "file": entry['file_name'],
                                    "status": "deduped",
                                })
                            continue

                    # Extract text
                    text = self._extract_text(entry['file_path'], entry['file_ext'])
                    if not text:
                        db.update_catalog_entry_status(entry['id'], 'error',
                                                       'No text extracted')
                        counts["errors"] += 1
                        if stream:
                            _stream_event({
                                "event": "progress",
                                "phase": "summarize",
                                "processed": processed_total,
                                "file": entry['file_name'],
                                "status": "error",
                                "error": "No text extracted",
                            })
                        continue

                    # Truncate for LLM
                    text_truncated = text[:8000]

                    # Call LLM for summary
                    result = self._llm_summarize(entry['file_name'], text_truncated)

                    # Update entry
                    db.update_catalog_entry_summary(
                        entry_id=entry['id'],
                        title=result.get('title', entry['file_name']),
                        summary=result.get('summary', ''),
                        tags=result.get('tags', ''),
                    )

                    # Embed summary into vector store. A failure here does not
                    # lose the summary, but it does leave the search index stale
                    # for this entry, so surface it as a warning rather than
                    # swallowing it.
                    embed_error = self._embed_summary(entry['id'], result)

                    counts["summarized"] += 1
                    if embed_error:
                        counts["embed_warnings"] = counts.get("embed_warnings", 0) + 1
                    if stream:
                        event = {
                            "event": "progress",
                            "phase": "summarize",
                            "processed": processed_total,
                            "file": entry['file_name'],
                            "status": "summarized",
                        }
                        if embed_error:
                            event["warning"] = f"embedding failed: {embed_error}"
                        _stream_event(event)

                except Exception as e:
                    db.update_catalog_entry_status(entry['id'], 'error', str(e))
                    counts["errors"] += 1
                    if stream:
                        _stream_event({
                            "event": "progress",
                            "phase": "summarize",
                            "processed": processed_total,
                            "file": entry['file_name'],
                            "status": "error",
                            "error": str(e),
                        })

        if stream:
            _stream_event({
                "event": "complete",
                "phase": "summarize",
                **counts,
            })

        return counts

    def _extract_text(self, file_path: str, file_ext: str) -> Optional[str]:
        """Extract text from a file using converters.

        Returns None only when the file type genuinely has no extractable text
        (an unsupported extension). A real extraction failure -- missing
        converter, unreadable file, corrupt document, parse error -- is allowed
        to propagate so the caller records the actual cause instead of a generic
        "No text extracted". Do not wrap this in a catch-all that hides the error.
        """
        try:
            from .converters import convert_to_markdown, is_supported
        except ImportError:
            from converters import convert_to_markdown, is_supported

        path = Path(file_path)

        # Direct text-based files
        text_extensions = {'.txt', '.md', '.csv', '.sql', '.json', '.xml'}
        if file_ext in text_extensions:
            return path.read_text(encoding='utf-8', errors='replace')

        # HTML files
        if file_ext in ('.html', '.htm'):
            return self._extract_html(path)

        # PowerPoint
        if file_ext == '.pptx':
            return self._extract_pptx(path)

        # Excel
        if file_ext == '.xlsx':
            return self._extract_xlsx(path)

        # Use existing converters for pdf, docx
        if is_supported(path):
            content, _meta = convert_to_markdown(path)
            return content

        # Genuinely unsupported file type -- no text to extract.
        return None

    def _extract_html(self, path: Path) -> str:
        """Extract text from HTML using built-in html.parser."""
        from html.parser import HTMLParser

        class TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.parts: List[str] = []
                self._skip = False

            def handle_starttag(self, tag, attrs):
                if tag in ('script', 'style'):
                    self._skip = True

            def handle_endtag(self, tag):
                if tag in ('script', 'style'):
                    self._skip = False

            def handle_data(self, data):
                if not self._skip:
                    text = data.strip()
                    if text:
                        self.parts.append(text)

        raw = path.read_text(encoding='utf-8', errors='replace')
        extractor = TextExtractor()
        extractor.feed(raw)
        return '\n'.join(extractor.parts)

    def _extract_pptx(self, path: Path) -> str:
        """Extract text from PowerPoint using python-pptx."""
        from pptx import Presentation

        prs = Presentation(str(path))
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)
            if slide_text:
                parts.append(f"## Slide {slide_num}\n" + '\n'.join(slide_text))
        return '\n\n'.join(parts)

    def _extract_xlsx(self, path: Path) -> str:
        """Extract text from Excel using openpyxl."""
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = []
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                if row_count >= 100:
                    rows_text.append("... (truncated at 100 rows)")
                    break
                cells = [str(c) if c is not None else '' for c in row]
                if any(cells):
                    rows_text.append(' | '.join(cells))
                row_count += 1
            if rows_text:
                parts.append(f"## Sheet: {sheet_name}\n" + '\n'.join(rows_text))
        wb.close()
        return '\n\n'.join(parts)

    def _llm_summarize(self, file_name: str, text: str) -> Dict[str, str]:
        """Call LLM to generate title, summary, and tags."""
        import openai

        try:
            from .config import OPENAI_API_KEY
        except ImportError:
            from config import OPENAI_API_KEY

        if not OPENAI_API_KEY:
            raise RuntimeError("OPENAI_API_KEY not set")

        client = openai.OpenAI(api_key=OPENAI_API_KEY)

        prompt = (
            f"Analyze this document (filename: {file_name}) and provide:\n"
            f"1. A concise title (max 80 chars)\n"
            f"2. A 2-3 sentence summary\n"
            f"3. 3-5 tags (comma-separated)\n\n"
            f"Respond in JSON format: "
            f'{{"title": "...", "summary": "...", "tags": "tag1, tag2, tag3"}}\n\n'
            f"Document text:\n{text}"
        )

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"},
            temperature=0.3,
            max_tokens=500,
        )

        result_text = response.choices[0].message.content
        return json.loads(result_text)

    def _embed_summary(self, entry_id: int, result: Dict[str, str]) -> Optional[str]:
        """Embed the summary into the vector store.

        Embedding failure is non-fatal to the summarize pass (the summary itself
        is already stored), but it must not be hidden: the search index is now
        out of date for this entry. Returns None on success, or a short
        description of the failure (exception type and message) so the caller can
        surface a warning instead of silently dropping it. The failure is
        returned, not raised, so a stale search index for one entry does not
        abort the whole summarize pass (the summary itself is already persisted).
        """
        try:
            try:
                from .vectors import embed_and_store
            except ImportError:
                from vectors import embed_and_store

            text_to_embed = f"{result.get('title', '')} {result.get('summary', '')}"
            embed_and_store(
                doc_id=f"catalog_{entry_id}",
                text=text_to_embed,
                collection="catalog",
                metadata=json.dumps({
                    "entry_id": entry_id,
                    "tags": result.get('tags', ''),
                }),
            )
        except Exception as exc:
            return f"{type(exc).__name__}: {exc}"
        return None
